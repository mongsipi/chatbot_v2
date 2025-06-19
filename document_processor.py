# document_processor.py - 안정화된 다중 문서 형식 처리 및 벡터 검색 모듈
import os
import json
import pickle
import re
from datetime import datetime
from typing import List, Dict, Tuple, Optional
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity

# PDF 처리
import PyPDF2
import pdfplumber

# Office 문서 처리 (안전한 로딩)
try:
    from docx import Document
    HAS_DOCX = True
    print("✓ python-docx 로딩 성공")
except ImportError as e:
    HAS_DOCX = False
    print(f"⚠️ python-docx 없음: Word 파일 지원 안함 - {e}")

try:
    from pptx import Presentation
    HAS_PPTX = True
    print("✓ python-pptx 로딩 성공")
except ImportError as e:
    HAS_PPTX = False
    print(f"⚠️ python-pptx 없음: PowerPoint 파일 지원 안함 - {e}")

try:
    import openpyxl
    import pandas as pd
    HAS_EXCEL = True
    print("✓ Excel 라이브러리 로딩 성공")
except ImportError as e:
    HAS_EXCEL = False
    print(f"⚠️ openpyxl/pandas 없음: Excel 파일 지원 안함 - {e}")

# sentence-transformers (안전한 로딩)
try:
    from sentence_transformers import SentenceTransformer
    HAS_SENTENCE_TRANSFORMERS = True
    print("✓ sentence-transformers 로딩 성공")
except ImportError as e:
    HAS_SENTENCE_TRANSFORMERS = False
    print(f"⚠️ sentence-transformers 없음: 기본 임베딩 사용 - {e}")

class SimpleEmbedding:
    """sentence-transformers가 없을 때 사용하는 간단한 임베딩"""
    
    def __init__(self):
        self.vocabulary = {}
        self.vocab_size = 0
        print("SimpleEmbedding 초기화 (fallback 모드)")
    
    def _build_vocabulary(self, texts):
        """텍스트에서 어휘 구축"""
        words = set()
        for text in texts:
            if isinstance(text, str):
                words.update(text.lower().split())
        
        self.vocabulary = {word: idx for idx, word in enumerate(sorted(words))}
        self.vocab_size = len(self.vocabulary)
    
    def encode(self, texts):
        """텍스트를 벡터로 변환"""
        if isinstance(texts, str):
            texts = [texts]
        
        if not self.vocabulary:
            self._build_vocabulary(texts)
        
        vectors = []
        for text in texts:
            vector = np.zeros(self.vocab_size)
            if isinstance(text, str):
                words = text.lower().split()
                for word in words:
                    if word in self.vocabulary:
                        vector[self.vocabulary[word]] += 1
            
            # 정규화
            if np.linalg.norm(vector) > 0:
                vector = vector / np.linalg.norm(vector)
            
            vectors.append(vector)
        
        return np.array(vectors)

class DocumentProcessor:
    """안정화된 다중 문서 형식 처리 및 벡터 검색 클래스"""
    
    def __init__(self, upload_folder):
        print(f"\n=== DocumentProcessor 초기화 시작 ===")
        
        self.upload_folder = upload_folder
        self.embeddings_file = os.path.join(upload_folder, 'embeddings.pkl')
        self.metadata_file = os.path.join(upload_folder, 'metadata.json')
        
        # 실제 지원 가능한 파일 확장자만 포함
        self.supported_extensions = {'.pdf': 'PDF'}
        
        if HAS_DOCX:
            self.supported_extensions['.docx'] = 'Word'
        if HAS_PPTX:
            self.supported_extensions['.pptx'] = 'PowerPoint'
        if HAS_EXCEL:
            self.supported_extensions['.xlsx'] = 'Excel'
            self.supported_extensions['.xls'] = 'Excel (Legacy)'
        
        print(f"지원 파일 형식: {list(self.supported_extensions.values())}")
        
        # 임베딩 모델 안전 초기화
        self.encoder = None
        self._safe_init_encoder()
        
        # 데이터 저장소
        self.documents = []
        self.embeddings = None
        self.metadata = {}
        
        # 데이터 로드
        self.load_data()
        
        print(f"=== DocumentProcessor 초기화 완료 ===")
        print(f"로드된 문서: {len(self.documents)}개")
    
    def _safe_init_encoder(self):
        """임베딩 모델 안전 초기화"""
        if HAS_SENTENCE_TRANSFORMERS:
            try:
                print("sentence-transformers 모델 로딩 시도...")
                # 가장 안정적인 모델 순서대로 시도
                models_to_try = [
                    'all-MiniLM-L6-v2',
                    'paraphrase-MiniLM-L6-v2', 
                    'all-mpnet-base-v2'
                ]
                
                for model_name in models_to_try:
                    try:
                        print(f"모델 시도: {model_name}")
                        self.encoder = SentenceTransformer(model_name)
                        print(f"✓ {model_name} 로딩 성공")
                        return
                    except Exception as e:
                        print(f"✗ {model_name} 로딩 실패: {e}")
                        continue
                
                print("모든 sentence-transformers 모델 로딩 실패, SimpleEmbedding 사용")
                self.encoder = SimpleEmbedding()
                
            except Exception as e:
                print(f"sentence-transformers 초기화 전체 실패: {e}")
                self.encoder = SimpleEmbedding()
        else:
            print("sentence-transformers 없음, SimpleEmbedding 사용")
            self.encoder = SimpleEmbedding()
    
    def extract_text_from_pdf(self, pdf_path):
        """PDF에서 텍스트 추출 (안정화)"""
        print(f"PDF 텍스트 추출: {os.path.basename(pdf_path)}")
        text = ""
        
        try:
            # pdfplumber 우선 시도
            with pdfplumber.open(pdf_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text += page_text + "\n"
                    except Exception as e:
                        print(f"페이지 {i+1} 처리 오류: {e}")
                        continue
                        
        except Exception as e1:
            print(f"pdfplumber 실패, PyPDF2 시도: {e1}")
            
            try:
                with open(pdf_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for i, page in enumerate(pdf_reader.pages):
                        try:
                            page_text = page.extract_text()
                            if page_text and page_text.strip():
                                text += page_text + "\n"
                        except Exception as e:
                            print(f"PyPDF2 페이지 {i+1} 처리 오류: {e}")
                            continue
            except Exception as e2:
                print(f"PyPDF2도 실패: {e2}")
                return None
        
        extracted_text = text.strip() if text.strip() else None
        if extracted_text:
            print(f"✓ PDF 텍스트 추출 성공: {len(extracted_text)}자")
        else:
            print("✗ PDF 텍스트 추출 실패")
        
        return extracted_text
    
    def extract_text_from_docx(self, docx_path):
        """Word 문서에서 텍스트 추출 (안전화)"""
        if not HAS_DOCX:
            return None
        
        try:
            print(f"Word 문서 처리: {os.path.basename(docx_path)}")
            doc = Document(docx_path)
            text = ""
            
            # 단락별 텍스트 추출
            for paragraph in doc.paragraphs:
                if paragraph.text and paragraph.text.strip():
                    text += paragraph.text.strip() + "\n"
            
            # 표 텍스트 추출 (안전하게)
            try:
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text and cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text += " | ".join(row_text) + "\n"
            except Exception as e:
                print(f"Word 표 처리 오류 (무시): {e}")
            
            result = text.strip() if text.strip() else None
            if result:
                print(f"✓ Word 텍스트 추출 성공: {len(result)}자")
            return result
            
        except Exception as e:
            print(f"✗ Word 문서 처리 오류: {e}")
            return None
    
    def extract_text_from_pptx(self, pptx_path):
        """PowerPoint 문서에서 텍스트 추출 (안전화)"""
        if not HAS_PPTX:
            return None
        
        try:
            print(f"PowerPoint 문서 처리: {os.path.basename(pptx_path)}")
            prs = Presentation(pptx_path)
            text = ""
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"\n=== 슬라이드 {i+1} ===\n"
                
                # 텍스트 상자에서 텍스트 추출
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            slide_text += shape.text.strip() + "\n"
                        
                        # 표 처리 (안전하게)
                        if hasattr(shape, 'has_table') and shape.has_table:
                            try:
                                table = shape.table
                                for row in table.rows:
                                    row_text = []
                                    for cell in row.cells:
                                        if cell.text and cell.text.strip():
                                            row_text.append(cell.text.strip())
                                    if row_text:
                                        slide_text += " | ".join(row_text) + "\n"
                            except Exception as e:
                                print(f"PowerPoint 표 처리 오류 (무시): {e}")
                                
                    except Exception as e:
                        print(f"PowerPoint shape 처리 오류 (무시): {e}")
                        continue
                
                if slide_text.strip() != f"=== 슬라이드 {i+1} ===":
                    text += slide_text
            
            result = text.strip() if text.strip() else None
            if result:
                print(f"✓ PowerPoint 텍스트 추출 성공: {len(result)}자")
            return result
            
        except Exception as e:
            print(f"✗ PowerPoint 문서 처리 오류: {e}")
            return None
    
    def extract_text_from_excel(self, excel_path):
        """Excel 문서에서 텍스트 추출 (안전화)"""
        if not HAS_EXCEL:
            return None
        
        try:
            print(f"Excel 문서 처리: {os.path.basename(excel_path)}")
            text = ""
            
            # pandas로 안전하게 읽기
            excel_file = pd.ExcelFile(excel_path)
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_path, sheet_name=sheet_name)
                    
                    if not df.empty:
                        sheet_text = f"\n=== {sheet_name} 시트 ===\n"
                        
                        # 컬럼명 추가
                        if not df.columns.empty:
                            valid_columns = [str(col) for col in df.columns if str(col).strip()]
                            if valid_columns:
                                sheet_text += "컬럼: " + " | ".join(valid_columns) + "\n"
                        
                        # 데이터 추가 (텍스트만)
                        for index, row in df.iterrows():
                            row_text = []
                            for value in row:
                                if pd.notna(value) and str(value).strip():
                                    str_value = str(value).strip()
                                    # 의미있는 텍스트만 포함 (숫자만 있는 것 제외)
                                    if len(str_value) > 1 and not str_value.replace('.', '').replace('-', '').replace(',', '').isdigit():
                                        row_text.append(str_value)
                            
                            if row_text:
                                sheet_text += " | ".join(row_text) + "\n"
                        
                        text += sheet_text
                        
                except Exception as e:
                    print(f"Excel 시트 {sheet_name} 처리 오류 (무시): {e}")
                    continue
            
            result = text.strip() if text.strip() else None
            if result:
                print(f"✓ Excel 텍스트 추출 성공: {len(result)}자")
            return result
            
        except Exception as e:
            print(f"✗ Excel 문서 처리 오류: {e}")
            return None
    
    def extract_text_from_document(self, file_path):
        """파일 형식에 따라 적절한 텍스트 추출 방법 선택"""
        filename = os.path.basename(file_path)
        _, ext = os.path.splitext(filename.lower())
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self.extract_text_from_excel(file_path)
        else:
            print(f"지원하지 않는 파일 형식: {ext}")
            return None
    
    def clean_text(self, text, file_type=None):
        """텍스트 정리 (간소화 및 안전화)"""
        if not text or not isinstance(text, str):
            return ""
        
        try:
            # 기본 정리만 수행 (복잡한 정규식 최소화)
            # 1. 제어문자 제거
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
            
            # 2. 기본 특수문자 정규화
            text = text.replace('·', ' ')
            text = text.replace('ㆍ', ' ')
            
            # 3. 과도한 공백 정리
            text = re.sub(r'\s+', ' ', text)
            text = re.sub(r'\n\s*\n+', '\n', text)
            
            # 4. 기본 줄 필터링
            lines = []
            for line in text.split('\n'):
                line = line.strip()
                if line and len(line) > 2:  # 너무 짧은 줄 제외
                    lines.append(line)
            
            cleaned_text = '\n'.join(lines)
            return cleaned_text
            
        except Exception as e:
            print(f"텍스트 정리 오류: {e}")
            return text  # 오류 시 원본 반환
    
    def chunk_text(self, text, chunk_size=600, overlap=50):
        """텍스트를 청크로 분할 (안전화)"""
        if not text:
            return []
        
        try:
            # 간단한 청킹 방식 사용
            chunks = []
            
            # 문단별로 분할
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
            
            current_chunk = ""
            for paragraph in paragraphs:
                # 현재 청크에 추가 가능한지 확인
                if len(current_chunk + paragraph) <= chunk_size:
                    current_chunk += paragraph + "\n"
                else:
                    # 현재 청크 저장
                    if current_chunk.strip():
                        chunks.append(current_chunk.strip())
                    
                    # 새 청크 시작
                    if len(paragraph) > chunk_size:
                        # 너무 긴 문단은 강제 분할
                        words = paragraph.split()
                        temp_chunk = ""
                        for word in words:
                            if len(temp_chunk + word) <= chunk_size:
                                temp_chunk += word + " "
                            else:
                                if temp_chunk.strip():
                                    chunks.append(temp_chunk.strip())
                                temp_chunk = word + " "
                        current_chunk = temp_chunk
                    else:
                        current_chunk = paragraph + "\n"
            
            # 마지막 청크 추가
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            
            # 빈 청크 제거
            chunks = [chunk for chunk in chunks if chunk.strip()]
            
            return chunks if chunks else [text[:chunk_size]]
            
        except Exception as e:
            print(f"청킹 오류: {e}")
            return [text]  # 오류 시 전체 텍스트를 하나의 청크로
    
    def process_document(self, file_path):
        """문서 파일 처리 (안전화)"""
        try:
            filename = os.path.basename(file_path)
            _, ext = os.path.splitext(filename.lower())
            
            print(f"\n=== 문서 처리: {filename} ===")
            
            # 지원 형식 확인
            if ext not in self.supported_extensions:
                print(f"✗ 지원하지 않는 파일 형식: {ext}")
                return False
            
            file_type = self.supported_extensions[ext]
            
            # 텍스트 추출
            text = self.extract_text_from_document(file_path)
            if not text:
                print(f"✗ 텍스트 추출 실패: {filename}")
                return False
            
            # 텍스트 정리
            cleaned_text = self.clean_text(text, file_type)
            if len(cleaned_text) < 50:  # 너무 짧은 텍스트 제외
                print(f"✗ 텍스트가 너무 짧음: {filename} ({len(cleaned_text)}자)")
                return False
            
            # 청킹
            chunks = self.chunk_text(cleaned_text)
            if not chunks:
                print(f"✗ 청크 생성 실패: {filename}")
                return False
            
            # 기존 문서가 있다면 제거
            self.remove_document_by_filename(filename)
            
            # 문서 추가
            for i, chunk in enumerate(chunks):
                doc = {
                    'content': chunk,
                    'filename': filename,
                    'file_type': file_type,
                    'chunk_id': i,
                    'file_path': file_path,
                    'processed_date': datetime.now().isoformat()
                }
                self.documents.append(doc)
            
            # 임베딩 생성
            self.update_embeddings()
            
            # 메타데이터 업데이트
            self.metadata[filename] = {
                'file_path': file_path,
                'file_type': file_type,
                'chunks_count': len(chunks),
                'processed_date': datetime.now().isoformat(),
                'file_size': os.path.getsize(file_path)
            }
            
            # 데이터 저장
            self.save_data()
            
            print(f"✓ 처리 완료: {filename} ({file_type}, {len(chunks)} 청크)")
            return True
            
        except Exception as e:
            print(f"✗ 문서 처리 오류 {file_path}: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def update_embeddings(self):
        """문서들의 임베딩 업데이트 (안전화)"""
        if not self.documents:
            self.embeddings = None
            return
        
        try:
            contents = [doc['content'] for doc in self.documents if doc.get('content')]
            if contents and self.encoder:
                self.embeddings = self.encoder.encode(contents)
                print(f"✓ 임베딩 업데이트: {len(contents)} 문서")
            else:
                print("⚠️ 임베딩 생성 스킵 (내용 없음 또는 인코더 없음)")
                self.embeddings = None
        except Exception as e:
            print(f"⚠️ 임베딩 생성 오류: {e}")
            self.embeddings = None
    
    def search_similar_documents(self, query, top_k=5, min_similarity=0.01):
        """다중 검색 방법을 사용한 문서 검색 (안전화)"""
        if not self.documents:
            return []
        
        print(f"\n=== 검색: '{query}' ===")
        
        all_results = []
        
        # 1. 키워드 기반 검색 (가장 안정적)
        keyword_results = self.keyword_search(query, top_k)
        all_results.extend(keyword_results)
        
        # 2. 벡터 임베딩 검색 (있는 경우에만)
        if self.embeddings is not None:
            try:
                vector_results = self.vector_search(query, top_k, min_similarity)
                all_results.extend(vector_results)
            except Exception as e:
                print(f"벡터 검색 오류 (무시): {e}")
        
        # 3. 부분 문자열 검색
        substring_results = self.substring_search(query, top_k)
        all_results.extend(substring_results)
        
        # 4. 결과 통합 및 중복 제거
        final_results = self.merge_and_rank_results(all_results, query)
        
        print(f"검색 완료: {len(final_results)}개 결과")
        return final_results[:top_k]
    
    def keyword_search(self, query, top_k=5):
        """키워드 기반 검색 (핵심 기능)"""
        keywords = self.extract_keywords(query)
        results = []
        
        for doc in self.documents:
            score = 0
            content_lower = doc['content'].lower()
            
            for keyword in keywords:
                keyword_lower = keyword.lower()
                count = content_lower.count(keyword_lower)
                if count > 0:
                    score += count * len(keyword)
            
            if score > 0:
                results.append({
                    'document': doc,
                    'similarity': min(score / 100.0, 1.0),
                    'content': doc['content'],
                    'method': 'keyword'
                })
        
        results.sort(key=lambda x: x['similarity'], reverse=True)
        return results[:top_k]
    
    def vector_search(self, query, top_k=5, min_similarity=0.01):
        """벡터 임베딩 검색"""
        try:
            query_embedding = self.encoder.encode([query])
            similarities = cosine_similarity(query_embedding, self.embeddings)[0]
            top_indices = np.argsort(similarities)[::-1][:top_k]
            
            results = []
            for idx in top_indices:
                similarity = similarities[idx]
                if similarity >= min_similarity:
                    results.append({
                        'document': self.documents[idx],
                        'similarity': float(similarity),
                        'content': self.documents[idx]['content'],
                        'method': 'vector'
                    })
            
            return results
        except Exception as e:
            print(f"벡터 검색 오류: {e}")
            return []
    
    def substring_search(self, query, top_k=5):
        """부분 문자열 검색"""
        results = []
        query_lower = query.lower()
        
        for doc in self.documents:
            content_lower = doc['content'].lower()
            
            if query_lower in content_lower:
                similarity = len(query) / len(doc['content'])
                results.append({
                    'document': doc,
                    'similarity': min(similarity * 10, 1.0),
                    'content': doc['content'],
                    'method': 'substring'
                })
        
        results.sort(key=lambda x: x['similarity'], reverse=True)
        return results[:top_k]
    
    def extract_keywords(self, query):
        """쿼리에서 키워드 추출"""
        stopwords = ['은', '는', '이', '가', '을', '를', '에', '에서', '와', '과', '의', '로', '으로', 
                    '에게', '한테', '께', '부터', '까지', '보다', '처럼', '같이', '와', '과', 
                    '그리고', '또는', '하지만', '그러나', '무엇', '어떤', '어디', '언제', '왜', '어떻게',
                    '뭐', '뭔', '뭘', '뭐가', '뭐를', '입니다', '습니다', '다', '요', '죠', '요?', '까요?']
        
        words = re.findall(r'[가-힣]+|[a-zA-Z]+|\d+', query)
        keywords = [word for word in words if len(word) >= 2 and word not in stopwords]
        
        return keywords
    
    def merge_and_rank_results(self, all_results, query):
        """검색 결과 통합 및 순위 결정"""
        if not all_results:
            return []
        
        doc_scores = {}
        
        for result in all_results:
            doc_id = id(result['document'])
            
            if doc_id not in doc_scores:
                doc_scores[doc_id] = {
                    'document': result['document'],
                    'content': result['content'],
                    'scores': [],
                    'methods': []
                }
            
            doc_scores[doc_id]['scores'].append(result['similarity'])
            doc_scores[doc_id]['methods'].append(result['method'])
        
        final_results = []
        for doc_id, info in doc_scores.items():
            if len(info['scores']) > 1:
                final_score = max(info['scores']) * 0.7 + sum(info['scores']) / len(info['scores']) * 0.3
            else:
                final_score = info['scores'][0]
            
            if 'keyword' in info['methods']:
                final_score *= 1.5
            
            final_results.append({
                'document': info['document'],
                'content': info['content'],
                'similarity': min(final_score, 1.0),
                'methods': info['methods']
            })
        
        final_results.sort(key=lambda x: x['similarity'], reverse=True)
        return final_results
    
    # 나머지 메서드들 (원본과 동일하지만 예외 처리 강화)
    def remove_document_by_filename(self, filename):
        """특정 파일의 모든 문서 제거"""
        self.documents = [doc for doc in self.documents if doc['filename'] != filename]
        if filename in self.metadata:
            del self.metadata[filename]
    
    def save_data(self):
        """데이터 저장"""
        try:
            if self.embeddings is not None:
                with open(self.embeddings_file, 'wb') as f:
                    pickle.dump({
                        'embeddings': self.embeddings,
                        'documents': self.documents
                    }, f)
            
            with open(self.metadata_file, 'w', encoding='utf-8') as f:
                json.dump(self.metadata, f, ensure_ascii=False, indent=2)
                
        except Exception as e:
            print(f"데이터 저장 오류: {e}")
    
    def load_data(self):
        """데이터 로드"""
        try:
            if os.path.exists(self.embeddings_file):
                with open(self.embeddings_file, 'rb') as f:
                    data = pickle.load(f)
                    self.embeddings = data.get('embeddings')
                    self.documents = data.get('documents', [])
            
            if os.path.exists(self.metadata_file):
                with open(self.metadata_file, 'r', encoding='utf-8') as f:
                    self.metadata = json.load(f)
                    
        except Exception as e:
            print(f"데이터 로드 오류: {e}")
            self.documents = []
            self.embeddings = None
            self.metadata = {}
    
    def get_uploaded_files(self):
        """업로드된 문서 파일 목록"""
        try:
            files = []
            for filename in os.listdir(self.upload_folder):
                _, ext = os.path.splitext(filename.lower())
                if ext in self.supported_extensions:
                    filepath = os.path.join(self.upload_folder, filename)
                    files.append({
                        'filename': filename,
                        'file_type': self.supported_extensions[ext],
                        'size': os.path.getsize(filepath),
                        'modified': datetime.fromtimestamp(os.path.getmtime(filepath)).isoformat()
                    })
            return files
        except Exception as e:
            print(f"파일 목록 조회 오류: {e}")
            return []
    
    def get_processed_files_info(self):
        """처리된 파일 정보"""
        return [
            {
                'filename': filename,
                'file_type': info['file_type'],
                'chunks_count': info['chunks_count'],
                'processed_date': info['processed_date'],
                'file_size': info.get('file_size', 0)
            }
            for filename, info in self.metadata.items()
        ]
    
    def has_processed_documents(self):
        """처리된 문서가 있는지 확인"""
        return len(self.documents) > 0
    
    def delete_file(self, filename):
        """파일 삭제"""
        try:
            filepath = os.path.join(self.upload_folder, filename)
            if os.path.exists(filepath):
                os.remove(filepath)
            
            self.remove_document_by_filename(filename)
            self.update_embeddings()
            self.save_data()
            
            return True
        except Exception as e:
            print(f"파일 삭제 오류: {e}")
            return False
    
    def reprocess_all_documents(self):
        """모든 문서 파일 재처리"""
        try:
            self.documents = []
            self.embeddings = None
            self.metadata = {}
            
            success_count = 0
            document_files = []
            
            for filename in os.listdir(self.upload_folder):
                _, ext = os.path.splitext(filename.lower())
                if ext in self.supported_extensions:
                    document_files.append(filename)
            
            for filename in document_files:
                filepath = os.path.join(self.upload_folder, filename)
                if self.process_document(filepath):
                    success_count += 1
            
            return success_count == len(document_files)
            
        except Exception as e:
            print(f"재처리 오류: {e}")
            return False
    
    def initialize_existing_documents(self):
        """서버 시작시 기존 문서 파일들 처리"""
        try:
            for filename in os.listdir(self.upload_folder):
                _, ext = os.path.splitext(filename.lower())
                if ext in self.supported_extensions:
                    if filename not in self.metadata:
                        filepath = os.path.join(self.upload_folder, filename)
                        print(f"기존 파일 처리: {filename}")
                        self.process_document(filepath)
                        
        except Exception as e:
            print(f"기존 파일 초기화 오류: {e}")


class QuestionAnalyzer:
    """질문 분석 및 답변 생성 클래스 (안전화)"""
    
    def __init__(self, document_processor):
        self.document_processor = document_processor
        
        self.greeting_patterns = [
            '안녕', 'hi', 'hello', '안녕하세요', '처음', '시작'
        ]
        
        self.thanks_patterns = [
            '감사', '고마워', '고맙', 'thank', '도움'
        ]
    
    def is_greeting(self, question):
        """인사말인지 확인"""
        question_lower = question.lower()
        return any(pattern in question_lower for pattern in self.greeting_patterns)
    
    def is_thanks(self, question):
        """감사 인사인지 확인"""
        question_lower = question.lower()
        return any(pattern in question_lower for pattern in self.thanks_patterns)
    
    def generate_greeting_response(self):
        """인사말 응답 생성"""
        supported_formats = ", ".join(self.document_processor.supported_extensions.values())
        
        return f"""안녕하세요! 😊

저는 다중 문서 형식을 지원하는 챗봇입니다. 
업로드된 문서의 내용을 분석하여 정확한 답변을 제공합니다.

**📁 지원 파일 형식**:
• {supported_formats}

**💡 이용 방법**:
• 문서 관련 질문을 자유롭게 해주세요
• 구체적인 질문일수록 더 정확한 답변을 받을 수 있습니다
• 문서에 없는 내용은 "해당 정보를 찾을 수 없습니다"라고 안내됩니다

궁금한 점이 있으시면 언제든 질문해주세요!"""
    
    def generate_thanks_response(self):
        """감사 인사 응답 생성"""
        return """천만에요! 😊

도움이 되셨다니 기쁩니다. 
다른 궁금한 점이 있으시면 언제든 질문해주세요!

**💡 팁**: 
• 더 구체적인 질문을 하시면 더 정확한 답변을 받을 수 있어요"""
    
    def analyze_question(self, question):
        """질문 분석 및 답변 생성 (안전화)"""
        try:
            if not question or len(question.strip()) < 2:
                return "질문을 입력해주세요."
            
            question = question.strip()
            print(f"\n=== 질문 분석: {question} ===")
            
            # 인사말 처리
            if self.is_greeting(question):
                return self.generate_greeting_response()
            
            # 감사 인사 처리
            if self.is_thanks(question):
                return self.generate_thanks_response()
            
            # 문서 검색
            search_results = self.document_processor.search_similar_documents(question, top_k=5, min_similarity=0.05)
            
            if not search_results:
                return self.generate_no_result_response_enhanced(question)
            
            # 답변 생성
            return self.generate_answer(question, search_results)
            
        except Exception as e:
            print(f"질문 분석 오류: {e}")
            import traceback
            traceback.print_exc()
            return "처리 중 오류가 발생했습니다. 다시 시도해주세요."
    
    def generate_no_result_response_enhanced(self, question):
        """결과가 없을 때 향상된 응답"""
        keywords = question.split()
        keyword_results = []
        
        for keyword in keywords:
            if len(keyword) > 1:
                for i, doc in enumerate(self.document_processor.documents):
                    if keyword.lower() in doc['content'].lower():
                        keyword_results.append({
                            'keyword': keyword,
                            'file_type': doc.get('file_type', 'Unknown'),
                            'content': doc['content'][:200] + "..."
                        })
                        break
        
        response = f'''**📋 "{question}"에 대한 검색 결과**

벡터 검색에서는 관련 내용을 찾지 못했습니다.'''
        
        if keyword_results:
            response += f"\n\n**🔍 키워드 기반 검색 결과:**\n"
            for result in keyword_results[:3]:
                response += f"\n• **'{result['keyword']}'** 관련 ({result['file_type']}):\n{result['content']}\n"
        
        file_types = {}
        for doc in self.document_processor.documents:
            file_type = doc.get('file_type', 'Unknown')
            file_types[file_type] = file_types.get(file_type, 0) + 1
        
        file_stats = ", ".join([f"{ft}: {count}개" for ft, count in file_types.items()])
        
        response += f'''

**💡 검색 개선 제안:**
• 다른 키워드로 질문해보세요
• 더 구체적이거나 더 일반적인 질문을 시도해보세요

**📊 현재 상태:**
• 처리된 파일 수: {len(self.document_processor.get_processed_files_info())}개
• 문서 타입별: {file_stats}
• 총 문서 청크: {len(self.document_processor.documents)}개

궁금한 점이 있으시면 다시 질문해주세요!'''
        
        return response
    
    def generate_answer(self, question, search_results):
        """검색 결과를 바탕으로 답변 생성"""
        try:
            best_results = [r for r in search_results if r['similarity'] > 0.3]
            
            if not best_results:
                return self.generate_no_result_response_enhanced(question)
            
            answer = f"**📋 '{question}'에 대한 답변**\n\n"
            
            main_content = best_results[0]['content']
            answer += f"{main_content}\n\n"
            
            if len(best_results) > 1:
                answer += "**📚 관련 추가 정보:**\n\n"
                for i, result in enumerate(best_results[1:3], 1):
                    content = result['content']
                    file_type = result['document'].get('file_type', 'Unknown')
                    if len(content) > 200:
                        content = content[:200] + "..."
                    answer += f"{i}. ({file_type}) {content}\n\n"
            
            sources_info = []
            for result in best_results:
                filename = result['document']['filename']
                file_type = result['document'].get('file_type', 'Unknown')
                sources_info.append(f"{filename} ({file_type})")
            
            unique_sources = list(dict.fromkeys(sources_info))
            answer += f"**📖 출처**: {', '.join(unique_sources)}\n\n"
            
            avg_similarity = sum(r['similarity'] for r in best_results) / len(best_results)
            confidence = "높음" if avg_similarity > 0.7 else "보통" if avg_similarity > 0.5 else "낮음"
            answer += f"**🎯 답변 신뢰도**: {confidence} ({avg_similarity:.2f})"
            
            return answer
            
        except Exception as e:
            print(f"답변 생성 오류: {e}")
            return "답변을 생성하는 중 오류가 발생했습니다."